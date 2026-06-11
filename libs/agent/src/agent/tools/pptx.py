from __future__ import annotations

import re
import uuid
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

from agent.models import SlideOutline


def _outputs_dir() -> Path:
    import os
    import tempfile

    env = os.environ.get("AGENT_OUTPUTS_DIR")
    if env:
        path = Path(env)
    else:
        path = Path(tempfile.gettempdir()) / "agent_outputs"
    path.mkdir(parents=True, exist_ok=True)
    return path


def _safe_filename(title: str) -> str:
    slug = re.sub(r"[^a-zA-Z0-9]+", "_", title.strip().lower()).strip("_")
    return slug[:60] or "lesson"


def create_pptx(outline: SlideOutline, run_id: str | None = None) -> Path:
    """Build a .pptx from a validated slide outline."""
    rid = run_id or uuid.uuid4().hex[:12]
    out_dir = _outputs_dir()
    filename = f"{_safe_filename(outline.title)}_{rid}.pptx"
    path = out_dir / filename

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    title_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_layout)
    title_slide.shapes.title.text = outline.title
    if title_slide.placeholders[1].text_frame:
        title_slide.placeholders[1].text = "Generated lesson slides"

    bullet_layout = prs.slide_layouts[1]
    for slide_spec in outline.slides:
        slide = prs.slides.add_slide(bullet_layout)
        slide.shapes.title.text = slide_spec.title
        body = slide.placeholders[1].text_frame
        body.clear()
        for index, bullet in enumerate(slide_spec.bullets):
            if index == 0:
                p = body.paragraphs[0]
            else:
                p = body.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.size = Pt(20)
        if slide_spec.speaker_note:
            notes = slide.notes_slide.notes_text_frame
            notes.text = slide_spec.speaker_note

    prs.save(str(path))
    return path
